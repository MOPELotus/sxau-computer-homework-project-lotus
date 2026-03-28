from __future__ import annotations

import base64
import io
import json
import math
import mimetypes
import re
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any

import fitz
import pandas as pd
from docx import Document
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from .data import normalize_dataset
from .llm import SiliconFlowAPIError, SiliconFlowClient


TABULAR_EXTENSIONS = {".csv", ".xlsx", ".xls"}
TEXT_EXTENSIONS = {".txt", ".md"}
JSON_EXTENSIONS = {".json"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
DOC_EXTENSIONS = {".doc"}
PPTX_EXTENSIONS = {".pptx"}
PPT_EXTENSIONS = {".ppt"}
ZIP_VISION_EXTENSIONS = DOCX_EXTENSIONS | PPTX_EXTENSIONS
MAX_VISION_IMAGES = 6
MAX_PDF_PAGES = 4

CANONICAL_COLUMNS = [
    "event_id",
    "report_date",
    "source_platform",
    "country",
    "border_region",
    "disease",
    "host_species",
    "cases",
    "deaths",
    "culling",
    "neighboring_outbreaks_14d",
    "livestock_density_index",
    "trade_flow_index",
    "transport_access_index",
    "border_distance_km",
    "rainfall_anomaly",
    "temperature_anomaly",
    "narrative_text",
    "cross_border_alert",
]

COLUMN_ALIASES = {
    "event_id": ["event_id", "id", "事件编号", "编号", "事件id", "记录编号"],
    "report_date": ["report_date", "日期", "报告日期", "上报日期", "发布时间", "通报日期"],
    "source_platform": ["source_platform", "来源平台", "平台", "数据来源", "来源", "来源系统"],
    "country": ["country", "国家", "发生国", "国家地区", "地区国家"],
    "border_region": ["border_region", "边境区域", "地区", "州省", "省州", "边境省份", "行政区"],
    "disease": ["disease", "疫病", "病种", "疾病名称", "疫病名称"],
    "host_species": ["host_species", "宿主动物", "宿主", "动物种类", "畜种", "动物"],
    "cases": ["cases", "发病数", "病例数", "感染数", "确诊数"],
    "deaths": ["deaths", "死亡数", "死亡病例", "死亡量"],
    "culling": ["culling", "扑杀数", "扑杀量", "无害化处理数"],
    "neighboring_outbreaks_14d": ["neighboring_outbreaks_14d", "14日周边疫情数", "近14日周边疫情数", "周边疫情数", "邻近疫情数"],
    "livestock_density_index": ["livestock_density_index", "养殖密度指数", "畜牧密度指数", "饲养密度指数"],
    "trade_flow_index": ["trade_flow_index", "贸易流指数", "贸易指数", "活畜贸易指数"],
    "transport_access_index": ["transport_access_index", "交通可达性指数", "交通指数", "运输便利度"],
    "border_distance_km": ["border_distance_km", "距边境距离公里", "距边境距离km", "边境距离", "距离边境公里"],
    "rainfall_anomaly": ["rainfall_anomaly", "降雨异常", "降雨距平"],
    "temperature_anomaly": ["temperature_anomaly", "气温异常", "温度异常", "气温距平"],
    "narrative_text": ["narrative_text", "疫情通报文本", "通报文本", "疫情描述", "文本", "摘要", "备注"],
    "cross_border_alert": ["cross_border_alert", "跨境预警标签", "预警标签", "风险标签", "标签", "alert", "label"],
}

TEXT_FIELD_ALIASES = {
    "report_date": [r"(?:日期|报告日期|上报日期|通报日期)\s*[:：]\s*(.+)"],
    "source_platform": [r"(?:来源平台|来源|平台)\s*[:：]\s*(.+)"],
    "country": [r"(?:国家|发生国)\s*[:：]\s*(.+)"],
    "border_region": [r"(?:边境区域|地区|省份|州省)\s*[:：]\s*(.+)"],
    "disease": [r"(?:疫病|病种|疾病名称)\s*[:：]\s*(.+)"],
    "host_species": [r"(?:宿主动物|宿主|动物种类|畜种)\s*[:：]\s*(.+)"],
    "cases": [r"(?:发病数|病例数|感染数)\s*[:：]\s*(\d+(?:\.\d+)?)"],
    "deaths": [r"(?:死亡数|死亡病例)\s*[:：]\s*(\d+(?:\.\d+)?)"],
    "culling": [r"(?:扑杀数|扑杀量)\s*[:：]\s*(\d+(?:\.\d+)?)"],
}


def load_dataset_from_source(
    source_path: str | Path,
    client: SiliconFlowClient | None = None,
) -> tuple[pd.DataFrame, dict[str, Any]]:
    path = Path(source_path)
    if not path.exists():
        raise FileNotFoundError(f"Source path not found: {path}")

    input_files = [path] if path.is_file() else sorted(file for file in path.rglob("*") if file.is_file())
    processed_frames: list[pd.DataFrame] = []
    report_items: list[dict[str, Any]] = []

    for file_path in input_files:
        frame, item = _ingest_single_file(file_path, client)
        report_items.append(item)
        if frame is not None and not frame.empty:
            processed_frames.append(frame)

    if not processed_frames:
        raise ValueError("未在资料目录中识别到可用于建模的资料。")

    merged = pd.concat(processed_frames, ignore_index=True)
    normalized = normalize_dataset(merged)

    report = {
        "source_path": str(path),
        "processed_files": [item for item in report_items if item["status"] == "processed"],
        "skipped_files": [item for item in report_items if item["status"] != "processed"],
        "processed_file_count": sum(1 for item in report_items if item["status"] == "processed"),
        "skipped_file_count": sum(1 for item in report_items if item["status"] != "processed"),
        "record_count": int(len(normalized)),
    }
    return normalized, report


def _ingest_single_file(
    file_path: Path,
    client: SiliconFlowClient | None,
) -> tuple[pd.DataFrame | None, dict[str, Any]]:
    suffix = file_path.suffix.lower()

    if suffix in TABULAR_EXTENSIONS:
        raw = _read_tabular_file(file_path)
        frame = _standardize_tabular_frame(raw, file_path.name, client)
        return frame, _processed_item(file_path, "tabular", len(frame))

    if suffix in TEXT_EXTENSIONS:
        content = _read_text_file(file_path)
        frame = pd.DataFrame(_extract_records_from_text(content, file_path.name, client))
        return frame, _processed_item(file_path, "text", len(frame))

    if suffix in JSON_EXTENSIONS:
        content = _read_text_file(file_path)
        try:
            payload = json.loads(content)
        except json.JSONDecodeError:
            payload = None
        if isinstance(payload, list) and payload and isinstance(payload[0], dict):
            frame = _standardize_tabular_frame(pd.DataFrame(payload), file_path.name, client)
            return frame, _processed_item(file_path, "json-table", len(frame))
        frame = pd.DataFrame(_extract_records_from_text(content, file_path.name, client))
        return frame, _processed_item(file_path, "json-text", len(frame))

    if suffix in DOCX_EXTENSIONS:
        text = _read_docx_text(file_path)
        if text.strip():
            frame = pd.DataFrame(_extract_records_from_text(text, file_path.name, client))
            return frame, _processed_item(file_path, "docx-text", len(frame))
        if client is not None:
            image_data_urls = _extract_zip_media_data_urls(file_path, prefix="word/media/")
            if image_data_urls:
                frame = pd.DataFrame(_extract_records_from_visual(image_data_urls, file_path.name, client))
                return frame, _processed_item(file_path, "docx-vision", len(frame))
        return None, _skipped_item(file_path, "docx", "未提取到可用文本或图片内容。")

    if suffix in PPTX_EXTENSIONS:
        text = _read_pptx_text(file_path)
        if text.strip():
            frame = pd.DataFrame(_extract_records_from_text(text, file_path.name, client))
            return frame, _processed_item(file_path, "pptx-text", len(frame))
        if client is not None:
            image_data_urls = _extract_zip_media_data_urls(file_path, prefix="ppt/media/")
            if image_data_urls:
                frame = pd.DataFrame(_extract_records_from_visual(image_data_urls, file_path.name, client))
                return frame, _processed_item(file_path, "pptx-vision", len(frame))
        return None, _skipped_item(file_path, "pptx", "未提取到可用文本或图片内容。")

    if suffix in PDF_EXTENSIONS:
        text = _read_pdf_text(file_path)
        if len(text.strip()) >= 20:
            frame = pd.DataFrame(_extract_records_from_text(text, file_path.name, client))
            return frame, _processed_item(file_path, "pdf-text", len(frame))
        if client is not None:
            image_data_urls = _pdf_to_data_urls(file_path)
            if image_data_urls:
                frame = pd.DataFrame(_extract_records_from_visual(image_data_urls, file_path.name, client))
                return frame, _processed_item(file_path, "pdf-vision", len(frame))
        return None, _skipped_item(file_path, "pdf", "PDF 未提取到有效文本，且当前无法继续进行视觉解析。")

    if suffix in IMAGE_EXTENSIONS:
        if client is None:
            return None, _skipped_item(file_path, "image", "图片资料需要配置 SiliconFlow API Key 才能调用视觉输入。")
        image_data_urls = [_local_image_to_data_url(file_path)]
        frame = pd.DataFrame(_extract_records_from_visual(image_data_urls, file_path.name, client))
        return frame, _processed_item(file_path, "image-vision", len(frame))

    if suffix in DOC_EXTENSIONS | PPT_EXTENSIONS:
        with tempfile.TemporaryDirectory() as temp_dir:
            converted_pdf = _convert_legacy_office_to_pdf(file_path, Path(temp_dir))
            if converted_pdf is not None and converted_pdf.exists():
                frame, _ = _ingest_single_file(converted_pdf, client)
                if frame is not None:
                    return frame, {
                        "file": str(file_path),
                        "status": "processed",
                        "kind": "legacy-office-converted",
                        "rows": int(len(frame)),
                        "note": f"已先转换为 PDF 再解析：{converted_pdf.name}",
                    }
        return None, _skipped_item(file_path, "legacy-office", "未能自动转换旧版 Office 文件。")

    return None, _skipped_item(file_path, "other", f"暂不支持的文件类型: {suffix or '无扩展名'}")


def _processed_item(file_path: Path, kind: str, rows: int) -> dict[str, Any]:
    return {
        "file": str(file_path),
        "status": "processed",
        "kind": kind,
        "rows": int(rows),
    }


def _skipped_item(file_path: Path, kind: str, note: str) -> dict[str, Any]:
    return {
        "file": str(file_path),
        "status": "skipped",
        "kind": kind,
        "rows": 0,
        "note": note,
    }


def _read_tabular_file(file_path: Path) -> pd.DataFrame:
    suffix = file_path.suffix.lower()
    if suffix == ".csv":
        encodings = ["utf-8-sig", "utf-8", "gb18030", "gbk"]
        last_error: Exception | None = None
        for encoding in encodings:
            try:
                return pd.read_csv(file_path, encoding=encoding)
            except UnicodeDecodeError as exc:
                last_error = exc
                continue
        raise last_error or ValueError(f"无法读取 CSV 文件: {file_path}")
    return pd.read_excel(file_path)


def _read_text_file(file_path: Path) -> str:
    encodings = ["utf-8-sig", "utf-8", "gb18030", "gbk"]
    for encoding in encodings:
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return file_path.read_text(encoding="latin-1")


def _read_pdf_text(file_path: Path) -> str:
    try:
        reader = PdfReader(str(file_path))
    except Exception:
        return ""
    pages = []
    for page in reader.pages[:20]:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(part for part in pages if part).strip()


def _read_docx_text(file_path: Path) -> str:
    try:
        document = Document(file_path)
    except Exception:
        return ""
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(parts)


def _read_pptx_text(file_path: Path) -> str:
    try:
        presentation = Presentation(file_path)
    except Exception:
        return ""
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                chunks.append(str(shape.text).strip())
    return "\n".join(chunks)


def _extract_zip_media_data_urls(file_path: Path, prefix: str) -> list[str]:
    data_urls: list[str] = []
    try:
        with zipfile.ZipFile(file_path) as archive:
            for name in archive.namelist():
                if not name.startswith(prefix):
                    continue
                mime_type = mimetypes.guess_type(name)[0] or "image/png"
                data_urls.append(_bytes_to_data_url(archive.read(name), mime_type))
                if len(data_urls) >= MAX_VISION_IMAGES:
                    break
    except Exception:
        return []
    return data_urls


def _pdf_to_data_urls(file_path: Path) -> list[str]:
    data_urls: list[str] = []
    try:
        document = fitz.open(file_path)
    except Exception:
        return []
    with document:
        for page_index in range(min(len(document), MAX_PDF_PAGES)):
            page = document.load_page(page_index)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            data_urls.append(_bytes_to_data_url(pixmap.tobytes("png"), "image/png"))
    return data_urls


def _local_image_to_data_url(file_path: Path) -> str:
    with Image.open(file_path) as image:
        buffer = io.BytesIO()
        image.convert("RGB").save(buffer, format="PNG")
    return _bytes_to_data_url(buffer.getvalue(), "image/png")


def _bytes_to_data_url(payload: bytes, mime_type: str) -> str:
    encoded = base64.b64encode(payload).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def _convert_legacy_office_to_pdf(file_path: Path, output_dir: Path) -> Path | None:
    converted = _convert_with_soffice(file_path, output_dir)
    if converted is not None:
        return converted
    return _convert_with_office_com(file_path, output_dir)


def _convert_with_soffice(file_path: Path, output_dir: Path) -> Path | None:
    soffice = shutil.which("soffice")
    if not soffice:
        return None
    file_path = file_path.resolve()
    output_dir = output_dir.resolve()
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(file_path)],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        return None
    candidate = output_dir / f"{file_path.stem}.pdf"
    return candidate if candidate.exists() else None


def _convert_with_office_com(file_path: Path, output_dir: Path) -> Path | None:
    file_path = file_path.resolve()
    output_dir = output_dir.resolve()
    output_pdf = output_dir / f"{file_path.stem}.pdf"
    escaped_input = str(file_path).replace("'", "''")
    escaped_output = str(output_pdf).replace("'", "''")

    if file_path.suffix.lower() == ".doc":
        script = f"""
$ErrorActionPreference = 'Stop'
$word = New-Object -ComObject Word.Application
$word.Visible = $false
$doc = $word.Documents.Open('{escaped_input}', $false, $true)
$doc.ExportAsFixedFormat('{escaped_output}', 17)
$doc.Close()
$word.Quit()
""".strip()
    else:
        script = f"""
$ErrorActionPreference = 'Stop'
$ppt = New-Object -ComObject PowerPoint.Application
$presentation = $ppt.Presentations.Open('{escaped_input}', $true, $false, $false)
$presentation.SaveAs('{escaped_output}', 32)
$presentation.Close()
$ppt.Quit()
""".strip()

    try:
        subprocess.run(
            ["powershell", "-NoProfile", "-Command", script],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        return None
    return output_pdf if output_pdf.exists() else None


def _standardize_tabular_frame(
    raw_frame: pd.DataFrame,
    source_file: str,
    client: SiliconFlowClient | None,
) -> pd.DataFrame:
    mapping = _infer_column_mapping(raw_frame.columns)
    if client is not None:
        missing_critical = [
            key
            for key in ["report_date", "country", "border_region", "disease", "host_species", "narrative_text"]
            if key not in mapping
        ]
        if missing_critical:
            try:
                llm_mapping = _llm_map_columns(raw_frame, client)
                mapping.update({key: value for key, value in llm_mapping.items() if value})
            except SiliconFlowAPIError:
                pass

    records = []
    for row_index, (_, row) in enumerate(raw_frame.iterrows(), start=1):
        record = _build_default_record(source_file, row_index)
        for canonical in CANONICAL_COLUMNS:
            if canonical in mapping:
                record[canonical] = row.get(mapping[canonical])
        record["narrative_text"] = _coerce_text(record.get("narrative_text")) or _compose_row_summary(row)
        record["event_id"] = _coerce_text(record.get("event_id")) or f"{Path(source_file).stem}-{row_index:04d}"
        record["report_date"] = _coerce_date(record.get("report_date"))
        record["source_platform"] = _coerce_text(record.get("source_platform")) or "资料汇总"
        record["country"] = _coerce_text(record.get("country")) or "待识别国家"
        record["border_region"] = _coerce_text(record.get("border_region")) or "待识别区域"
        record["disease"] = _coerce_text(record.get("disease")) or "待识别疫病"
        record["host_species"] = _coerce_text(record.get("host_species")) or "待识别宿主"
        record["cases"] = _coerce_number(record.get("cases"))
        record["deaths"] = _coerce_number(record.get("deaths"))
        record["culling"] = _coerce_number(record.get("culling"))
        record["neighboring_outbreaks_14d"] = _coerce_number(record.get("neighboring_outbreaks_14d"))
        record["livestock_density_index"] = _coerce_number(record.get("livestock_density_index"))
        record["trade_flow_index"] = _coerce_number(record.get("trade_flow_index"))
        record["transport_access_index"] = _coerce_number(record.get("transport_access_index"))
        record["border_distance_km"] = _coerce_number(record.get("border_distance_km"))
        record["rainfall_anomaly"] = _coerce_number(record.get("rainfall_anomaly"), allow_negative=True)
        record["temperature_anomaly"] = _coerce_number(record.get("temperature_anomaly"), allow_negative=True)
        label_value = record.get("cross_border_alert")
        record["label_available"] = _has_meaningful_value(label_value)
        record["cross_border_alert"] = _coerce_label(label_value)
        records.append(record)

    return pd.DataFrame(records)


def _extract_records_from_text(
    text: str,
    source_file: str,
    client: SiliconFlowClient | None,
) -> list[dict[str, Any]]:
    if client is not None:
        try:
            payload, _ = client.chat_json(
                system_prompt=(
                    "你是动物疫病资料整理助手。"
                    "请从文本中抽取标准化疫情记录，并仅返回 JSON 对象。"
                ),
                user_prompt=f"""
请从以下资料中抽取 1 条或多条动物疫病记录，返回 JSON：
{{
  "records": [
    {{
      "event_id": "若原文没有则留空",
      "report_date": "YYYY-MM-DD，没有则留空",
      "source_platform": "来源平台，没有则留空",
      "country": "国家",
      "border_region": "边境区域",
      "disease": "疫病名称",
      "host_species": "宿主动物",
      "cases": 数值,
      "deaths": 数值,
      "culling": 数值,
      "neighboring_outbreaks_14d": 数值,
      "livestock_density_index": 数值,
      "trade_flow_index": 数值,
      "transport_access_index": 数值,
      "border_distance_km": 数值,
      "rainfall_anomaly": 数值,
      "temperature_anomaly": 数值,
      "narrative_text": "疫情描述",
      "cross_border_alert": 0或1或空
    }}
  ]
}}

要求：
1. 没有给出的数值统一填 0。
2. 没有明确标签时，cross_border_alert 置空，不要强行猜测。
3. narrative_text 使用中文概括，不超过120字。

资料内容如下：
{text[:12000]}
""".strip(),
                temperature=0.1,
                max_tokens=1600,
            )
            records = payload.get("records", [])
            if isinstance(records, list) and records:
                return _normalize_llm_records(records, source_file, fallback_text=text)
        except SiliconFlowAPIError:
            pass

    record = _build_default_record(source_file, 1)
    for field, patterns in TEXT_FIELD_ALIASES.items():
        for pattern in patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                record[field] = match.group(1).strip()
                break

    record["narrative_text"] = re.sub(r"\s+", " ", text).strip()[:120] or "未提取到详细描述"
    record["event_id"] = _coerce_text(record.get("event_id")) or f"{Path(source_file).stem}-0001"
    record["report_date"] = _coerce_date(record.get("report_date"))
    record["source_platform"] = _coerce_text(record.get("source_platform")) or "文本资料"
    record["country"] = _coerce_text(record.get("country")) or "待识别国家"
    record["border_region"] = _coerce_text(record.get("border_region")) or "待识别区域"
    record["disease"] = _coerce_text(record.get("disease")) or "待识别疫病"
    record["host_species"] = _coerce_text(record.get("host_species")) or "待识别宿主"
    record["cases"] = _coerce_number(record.get("cases"))
    record["deaths"] = _coerce_number(record.get("deaths"))
    record["culling"] = _coerce_number(record.get("culling"))
    record["cross_border_alert"] = 0
    record["label_available"] = False
    return [record]


def _extract_records_from_visual(
    image_data_urls: list[str],
    source_file: str,
    client: SiliconFlowClient,
    supporting_text: str = "",
) -> list[dict[str, Any]]:
    prompt_text = f"""
请根据图片中的内容抽取 1 条或多条动物疫病记录，返回 JSON：
{{
  "records": [
    {{
      "event_id": "若看不到则留空",
      "report_date": "YYYY-MM-DD，没有则留空",
      "source_platform": "来源平台，没有则留空",
      "country": "国家",
      "border_region": "边境区域",
      "disease": "疫病名称",
      "host_species": "宿主动物",
      "cases": 数值,
      "deaths": 数值,
      "culling": 数值,
      "neighboring_outbreaks_14d": 数值,
      "livestock_density_index": 数值,
      "trade_flow_index": 数值,
      "transport_access_index": 数值,
      "border_distance_km": 数值,
      "rainfall_anomaly": 数值,
      "temperature_anomaly": 数值,
      "narrative_text": "中文概括，不超过120字",
      "cross_border_alert": 0或1或空
    }}
  ]
}}

要求：
1. 只返回 JSON。
2. 看不清或没有的信息填空或填 0，不要臆造。
3. 若下方附带了辅助文字，请结合辅助文字一起判断。

辅助文字：
{supporting_text[:2000] if supporting_text else "无"}
""".strip()

    content_parts: list[dict[str, Any]] = [{"type": "text", "text": prompt_text}]
    for data_url in image_data_urls[:MAX_VISION_IMAGES]:
        content_parts.append(
            {
                "type": "image_url",
                "image_url": {
                    "url": data_url,
                    "detail": "high",
                },
            }
        )

    payload, _ = client.chat_json(
        system_prompt="你是动物疫病资料视觉分析助手。请阅读图片或页面并抽取结构化疫情记录。",
        user_prompt=content_parts,
        temperature=0.1,
        max_tokens=1800,
    )
    records = payload.get("records", [])
    if not isinstance(records, list) or not records:
        raise SiliconFlowAPIError("视觉输入未返回可用 records。")
    return _normalize_llm_records(records, source_file, fallback_text=supporting_text)


def _normalize_llm_records(
    records: list[dict[str, Any]],
    source_file: str,
    fallback_text: str,
) -> list[dict[str, Any]]:
    normalized_records = []
    for index, record in enumerate(records, start=1):
        current = _build_default_record(source_file, index)
        current.update(record)
        current["narrative_text"] = _coerce_text(current.get("narrative_text")) or fallback_text[:120]
        current["event_id"] = _coerce_text(current.get("event_id")) or f"{Path(source_file).stem}-{index:04d}"
        current["report_date"] = _coerce_date(current.get("report_date"))
        current["source_platform"] = _coerce_text(current.get("source_platform")) or "多模态资料"
        current["country"] = _coerce_text(current.get("country")) or "待识别国家"
        current["border_region"] = _coerce_text(current.get("border_region")) or "待识别区域"
        current["disease"] = _coerce_text(current.get("disease")) or "待识别疫病"
        current["host_species"] = _coerce_text(current.get("host_species")) or "待识别宿主"
        current["cases"] = _coerce_number(current.get("cases"))
        current["deaths"] = _coerce_number(current.get("deaths"))
        current["culling"] = _coerce_number(current.get("culling"))
        current["neighboring_outbreaks_14d"] = _coerce_number(current.get("neighboring_outbreaks_14d"))
        current["livestock_density_index"] = _coerce_number(current.get("livestock_density_index"))
        current["trade_flow_index"] = _coerce_number(current.get("trade_flow_index"))
        current["transport_access_index"] = _coerce_number(current.get("transport_access_index"))
        current["border_distance_km"] = _coerce_number(current.get("border_distance_km"))
        current["rainfall_anomaly"] = _coerce_number(current.get("rainfall_anomaly"), allow_negative=True)
        current["temperature_anomaly"] = _coerce_number(current.get("temperature_anomaly"), allow_negative=True)
        current["label_available"] = _has_meaningful_value(current.get("cross_border_alert"))
        current["cross_border_alert"] = _coerce_label(current.get("cross_border_alert"))
        normalized_records.append(current)
    return normalized_records


def _llm_map_columns(raw_frame: pd.DataFrame, client: SiliconFlowClient) -> dict[str, str]:
    sample_rows = raw_frame.head(3).fillna("").to_dict(orient="records")
    payload, _ = client.chat_json(
        system_prompt=(
            "你是表格字段映射助手。"
            "请把原始列名映射到标准疫情字段，只返回 JSON 对象。"
        ),
        user_prompt=f"""
请根据原始列名和样例行，将它们映射到以下标准字段：
{CANONICAL_COLUMNS}

原始列名：
{list(raw_frame.columns)}

样例行：
{json.dumps(sample_rows, ensure_ascii=False)}

返回格式：
{{
  "column_mapping": {{
    "event_id": "原始列名或空",
    "report_date": "原始列名或空",
    "source_platform": "原始列名或空",
    "country": "原始列名或空",
    "border_region": "原始列名或空",
    "disease": "原始列名或空",
    "host_species": "原始列名或空",
    "cases": "原始列名或空",
    "deaths": "原始列名或空",
    "culling": "原始列名或空",
    "neighboring_outbreaks_14d": "原始列名或空",
    "livestock_density_index": "原始列名或空",
    "trade_flow_index": "原始列名或空",
    "transport_access_index": "原始列名或空",
    "border_distance_km": "原始列名或空",
    "rainfall_anomaly": "原始列名或空",
    "temperature_anomaly": "原始列名或空",
    "narrative_text": "原始列名或空",
    "cross_border_alert": "原始列名或空"
  }}
}}
""".strip(),
        temperature=0.0,
        max_tokens=1200,
    )
    mapping = payload.get("column_mapping", {})
    return {
        canonical: original
        for canonical, original in mapping.items()
        if canonical in CANONICAL_COLUMNS and original in raw_frame.columns
    }


def _infer_column_mapping(columns: list[str] | pd.Index) -> dict[str, str]:
    normalized_aliases = {
        _normalize_name(alias): canonical
        for canonical, aliases in COLUMN_ALIASES.items()
        for alias in aliases
    }
    mapping: dict[str, str] = {}
    for original in columns:
        normalized = _normalize_name(str(original))
        canonical = normalized_aliases.get(normalized)
        if canonical and canonical not in mapping:
            mapping[canonical] = str(original)
    return mapping


def _normalize_name(value: str) -> str:
    return re.sub(r"[\s_\-()（）\[\]【】:/]+", "", value).lower()


def _compose_row_summary(row: pd.Series) -> str:
    pieces = []
    for column, value in row.items():
        if _has_meaningful_value(value):
            pieces.append(f"{column}:{value}")
    return "；".join(pieces)[:160]


def _build_default_record(source_file: str, row_index: int) -> dict[str, Any]:
    return {
        "event_id": f"{Path(source_file).stem}-{row_index:04d}",
        "report_date": pd.Timestamp.today().strftime("%Y-%m-%d"),
        "source_platform": "资料汇总",
        "source_file": source_file,
        "country": "待识别国家",
        "border_region": "待识别区域",
        "disease": "待识别疫病",
        "host_species": "待识别宿主",
        "cases": 0,
        "deaths": 0,
        "culling": 0,
        "neighboring_outbreaks_14d": 0,
        "livestock_density_index": 0,
        "trade_flow_index": 0,
        "transport_access_index": 0,
        "border_distance_km": 0,
        "rainfall_anomaly": 0,
        "temperature_anomaly": 0,
        "narrative_text": "",
        "cross_border_alert": 0,
        "label_available": False,
    }


def _coerce_text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, float) and math.isnan(value):
        return ""
    return str(value).strip()


def _coerce_date(value: Any) -> str:
    text = _coerce_text(value)
    if not text:
        return pd.Timestamp.today().strftime("%Y-%m-%d")
    timestamp = pd.to_datetime(text, errors="coerce")
    if pd.isna(timestamp):
        return pd.Timestamp.today().strftime("%Y-%m-%d")
    return timestamp.strftime("%Y-%m-%d")


def _coerce_number(value: Any, allow_negative: bool = False) -> float:
    if value is None:
        return 0
    if isinstance(value, str):
        value = value.replace(",", "").strip()
    try:
        number = float(value)
    except (TypeError, ValueError):
        return 0
    if not allow_negative and number < 0:
        return 0
    return number


def _coerce_label(value: Any) -> int:
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {"1", "是", "高", "高风险", "high", "true"}:
            return 1
        if normalized in {"0", "否", "低", "低风险", "low", "false"}:
            return 0
    try:
        return 1 if int(float(value)) > 0 else 0
    except (TypeError, ValueError):
        return 0


def _has_meaningful_value(value: Any) -> bool:
    if value is None:
        return False
    if isinstance(value, float) and math.isnan(value):
        return False
    return str(value).strip() != ""
