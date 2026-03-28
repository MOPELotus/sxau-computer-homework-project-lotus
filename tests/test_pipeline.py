from __future__ import annotations

import json
import shutil
import tempfile
import unittest
from pathlib import Path

import fitz
from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation

from disease_intel.pipeline import run_pipeline


class PipelineSmokeTest(unittest.TestCase):
    def test_folder_pipeline_runs_in_heuristic_mode(self) -> None:
        project_root = Path(__file__).resolve().parent.parent
        source_root = project_root / "materials" / "source"

        with tempfile.TemporaryDirectory() as source_dir, tempfile.TemporaryDirectory() as output_dir:
            temp_source = Path(source_dir)
            for file_path in source_root.iterdir():
                if file_path.is_file():
                    shutil.copy2(file_path, temp_source / file_path.name)

            self._create_docx_file(temp_source / "补充材料.docx")
            self._create_pptx_file(temp_source / "汇报提纲.pptx")
            self._create_pdf_file(temp_source / "扫描通报.pdf")
            self._create_image_file(temp_source / "疫病截图.png")

            summary = run_pipeline(
                source_path=temp_source,
                output_dir=output_dir,
                llm_mode="heuristic",
            )

            self.assertEqual(summary["llm_mode"], "heuristic")
            self.assertEqual(summary["model_mode"], "supervised")

            metrics_path = Path(output_dir) / "metrics.json"
            predictions_path = Path(output_dir) / "predictions.csv"
            standardized_path = Path(output_dir) / "standardized_dataset.csv"
            ingest_report_path = Path(output_dir) / "ingest_report.json"

            self.assertTrue(metrics_path.exists())
            self.assertTrue(predictions_path.exists())
            self.assertTrue(standardized_path.exists())
            self.assertTrue(ingest_report_path.exists())

            metrics = json.loads(metrics_path.read_text(encoding="utf-8"))
            ingest_report = json.loads(ingest_report_path.read_text(encoding="utf-8"))
            self.assertIn("f1", metrics)
            self.assertGreaterEqual(ingest_report["record_count"], 12)
            self.assertGreaterEqual(ingest_report["processed_file_count"], 5)
            self.assertGreaterEqual(ingest_report["skipped_file_count"], 1)

    @staticmethod
    def _create_docx_file(path: Path) -> None:
        document = Document()
        document.add_heading("边境疫情补充说明", level=1)
        document.add_paragraph("国家：缅甸")
        document.add_paragraph("边境区域：掸邦")
        document.add_paragraph("报告日期：2025-03-26")
        document.add_paragraph("疫病：禽流感")
        document.add_paragraph("宿主动物：家禽")
        document.add_paragraph("发病数：64")
        document.add_paragraph("死亡数：18")
        document.add_paragraph("扑杀数：220")
        document.add_paragraph("疫情描述：多个边境乡镇报告家禽异常死亡，市场交易链条仍在运行。")
        document.save(path)

    @staticmethod
    def _create_pptx_file(path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "边境疫病简报"
        slide.placeholders[1].text = (
            "国家：尼泊尔\n"
            "边境区域：蓝毗尼省\n"
            "报告日期：2025-03-29\n"
            "疫病：牛结节性皮肤病\n"
            "宿主动物：牛\n"
            "发病数：83\n"
            "死亡数：6\n"
            "扑杀数：90\n"
            "疫情描述：病例与活畜运输和集市交易有关。"
        )
        presentation.save(path)

    @staticmethod
    def _create_pdf_file(path: Path) -> None:
        document = fitz.open()
        page = document.new_page()
        page.insert_text(
            (72, 72),
            "country: Laos\nregion: Luang Namtha\ndate: 2025-03-30\ndisease: ASF\nhost: swine\ncases: 51\ndeaths: 9\nculling: 150\nnarrative: farms near the border trade corridor reported unusual pig deaths.",
            fontsize=12,
        )
        document.save(path)
        document.close()

    @staticmethod
    def _create_image_file(path: Path) -> None:
        image = Image.new("RGB", (800, 400), color="white")
        draw = ImageDraw.Draw(image)
        draw.text((20, 20), "country: Vietnam disease: FMD visual-input smoke file", fill="black")
        image.save(path)


if __name__ == "__main__":
    unittest.main()
